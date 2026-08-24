import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Deep Reinforcement Learning for Proactive Network Traffic Congestion Control"
    slide.placeholders[1].text = "Authors: Peter Borngreat-Mensah, Henry Asante, Blessing Listowell Issaka\nMSc Computer Science, Cohort B, University of Ghana"
    slide.notes_slide.notes_text_frame.text = "Welcome the audience, supervisor, and committee. Introduce the research topic: evaluating Deep Reinforcement Learning (DQN and PPO) against static ECMP for proactive load balancing across wide-area ISP networks."
    
    # --- Slide 2: Problem Statement ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Problem: Static Routing & Micro-Bursts"
    slide.placeholders[1].text = (
        "- Wide-area networks suffer from transient micro-bursts and volatile traffic demands.\n"
        "- Equal-Cost Multi-Path (ECMP, RFC 2992) relies on static packet header hashing.\n"
        "- ECMP is stateless and blind to real-time link queue depths.\n"
        "- Hash collisions map heavy 'elephant flows' onto the same bottleneck links, causing packet drops."
    )
    slide.notes_slide.notes_text_frame.text = "Explain that traditional ECMP relies on static hashing of packet headers without inspecting buffer occupancy. When heavy elephant flows hash to the same link, severe queue overflows occur."

    # --- Slide 3: Research Objectives ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Research Objectives"
    slide.placeholders[1].text = (
        "1. Implement DRL agents (DQN and PPO) inside an SDN-controlled framework.\n"
        "2. Evaluate zero-shot generalization on 5 completely unseen ISP topologies.\n"
        "3. Evaluate per-step routing efficiency vs cumulative episode truncation effects.\n"
        "4. Identify theoretical capacity boundaries under heavy network saturation."
    )
    slide.notes_slide.notes_text_frame.text = "State our key objectives: evaluating DQN and PPO against ECMP, validating zero-shot transferability on held-out topologies, examining per-step dynamics, and identifying saturation limits."

    # --- Slide 4: Methodology - State Space & Telemetry ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Methodology: State Space & PCA Variance"
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4.5))
    tf = txBox.text_frame
    tf.text = "MDP State Formulation:"
    p = tf.add_paragraph()
    p.text = "- 250-dimensional continuous telemetry vector.\n- Top 100 links (utilization ratio).\n- Top 100 switch queues (buffer depth).\n- 50 instantaneous flow demands.\n- Deterministic zero-padding for topologies with < 100 links."
    p = tf.add_paragraph()
    p.text = "- PCA Variance: High variance in link and queue states motivated feature normalization."
    try:
        slide.shapes.add_picture('results/plots/feature_importance.png', Inches(5.0), Inches(1.5), width=Inches(4.5))
    except Exception as e:
        print(f"Skipped image: {e}")
    slide.notes_slide.notes_text_frame.text = "Explain the 250-feature state vector. Zero-padding allows the network to process any topology graph. PCA justified normalising queue and link inputs prior to feeding into the policy network."

    # --- Slide 5: Methodology - Action Space Dichotomy ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Methodology: Continuous vs Discrete Actions"
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4.5))
    tf = txBox.text_frame
    tf.text = "Action Space Architecture:"
    p = tf.add_paragraph()
    p.text = "- PPO: Explores continuous probability distributions over candidate K-shortest paths.\n- DQN: Bound to 10 discrete routing bins.\n- Structural Caveat: Comparing discrete DQN to continuous PPO is an intentional architectural contrast.\n- Reward Function: Penalizes queue congestion and packet drops while encouraging throughput."
    try:
        slide.shapes.add_picture('results/plots/action_space_dist.png', Inches(5.0), Inches(1.5), width=Inches(4.5))
    except Exception as e:
        pass
    slide.notes_slide.notes_text_frame.text = "Highlight the architectural dichotomy: PPO optimizes continuous load weights while DQN operates on 10 discrete bins. Note our structural caveat regarding the different action representations."

    # --- Slide 6: Moderate Traffic Performance ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Primary Success: Moderate Traffic Regime"
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4.5))
    tf = txBox.text_frame
    tf.text = "Optimal Operating Window:"
    p = tf.add_paragraph()
    p.text = "- Moderate traffic is where BOTH DRL agents cleanly outperform ECMP.\n- Zero-shot test on 5 held-out topologies (paired two-tailed t-test, df = 4):\n  * DQN vs ECMP: t = 151.32, p < 0.0001\n  * PPO vs ECMP: t = 76.49, p < 0.0001\n- DQN maintained lower median link utilization (0.88 vs 0.92), proactively avoiding bottleneck formation."
    try:
        slide.shapes.add_picture('results/plots/utilization_boxplot.png', Inches(5.0), Inches(1.5), width=Inches(4.5))
    except Exception as e:
        pass
    slide.notes_slide.notes_text_frame.text = "Point to the boxplot. Moderate traffic is the sweet spot. Both DQN and PPO demonstrate statistically significant superiority over ECMP across all held-out test networks."

    # --- Slide 7: Burst Traffic & Per-Step Efficiency ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Burst Traffic & Per-Step Efficiency Dynamics"
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4.5))
    tf = txBox.text_frame
    tf.text = "Per-Step vs Cumulative Reward:"
    p = tf.add_paragraph()
    p.text = "- Cumulative Reward: DQN (-49.48) beat ECMP (-54.12), while PPO scored -58.52.\n- Survival Length: PPO survived 1760 steps vs DQN's 1380 and ECMP's 1050 steps.\n- True Per-Step Reward:\n  * PPO: -0.0333 / step (Highest efficiency)\n  * DQN: -0.0359 / step\n  * ECMP: -0.0515 / step (Worst)\n- Note: Cumulative and per-step metrics are mechanically coupled through episode duration under fixed drop-threshold termination."
    try:
        slide.shapes.add_picture('results/plots/burst_performance_bar.png', Inches(5.0), Inches(1.5), width=Inches(4.5))
    except Exception as e:
        pass
    slide.notes_slide.notes_text_frame.text = "Explain the crucial mathematical nuance: PPO achieves the highest per-step reward (-0.0333 vs DQN's -0.0359 and ECMP's -0.0515) and survives 1760 steps. Episodic sums penalize it simply because it operates longer before truncation."

    # --- Slide 8: The Saturation Limit (Honest Negative Result) ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "The Saturation Limit (Honest Negative Result)"
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4.5))
    tf = txBox.text_frame
    tf.text = "Physical Capacity Boundaries:"
    p = tf.add_paragraph()
    p.text = "- Under heavy congestion, all links reach 100% saturation.\n- DRL performed worse than ECMP (-80.41 and -80.43 vs -73.78).\n- Continued path-shifting under saturation creates counterproductive routing churn.\n- Proves that routing algorithms cannot create bandwidth without admission control."
    try:
        slide.shapes.add_picture('results/plots/throughput_vs_latency.png', Inches(5.0), Inches(1.5), width=Inches(4.5))
    except Exception as e:
        pass
    slide.notes_slide.notes_text_frame.text = "Discuss the honest negative result: When the network is globally saturated, no routing policy can circumvent bottlenecks. Path exploration under full saturation causes extra churn."

    # --- Slide 9: Zero-Shot Generalization Across Topologies ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Zero-Shot Generalization on 5 Unseen Topologies"
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4.5))
    tf = txBox.text_frame
    tf.text = "Held-Out Test Set (Zero Overlap):"
    p = tf.add_paragraph()
    p.text = "- Evaluated on Abilene (11 nodes), Janetbackbone (29 nodes), Renater2001 (24 nodes), Bellcanada (48 nodes), and Colt (153 nodes).\n- DQN achieved uniform improvement (+6.85 reward units) across all scales.\n- Uniformity enabled by 250-D state vector zero-padding.\n- Single-palette heatmap illustrates consistent performance gradients."
    try:
        slide.shapes.add_picture('results/plots/reward_heatmap.png', Inches(5.0), Inches(1.5), width=Inches(4.5))
    except Exception as e:
        pass
    slide.notes_slide.notes_text_frame.text = "Point to the reward matrix. Our zero-shot evaluation on 5 completely unseen topologies validates that the learned routing policies transfer successfully across unfamiliar graph structures spanning 11 to 153 nodes."

    # --- Slide 10: Conclusion & Recommendations ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusions & Recommendations"
    slide.placeholders[1].text = (
        "- DRL is highly effective for proactive load balancing under Moderate traffic (both win) and Burst traffic (DQN wins cumulative, PPO wins per-step).\n"
        "- Per-step analysis revealed PPO's superior efficiency (-0.0333/step) and prolonged survival (1760 steps).\n"
        "- Zero-shot topological generalization is viable via standardized telemetry zero-padding.\n"
        "- DRL routing cannot solve physical saturation without active admission control.\n"
        "- Recommendation: Combine fixed-horizon benchmarking with dynamic admission control."
    )
    slide.notes_slide.notes_text_frame.text = "Conclude the presentation by summarizing core takeaways: DRL excels under moderate and burst conditions, per-step metrics resolve cumulative truncation artifacts, and admission control is essential for saturation. Thank the committee and invite questions."

    prs.save('final_presentation.pptx')
    print("final_presentation.pptx successfully created with 10 slides, presenter notes, and updated visuals.")

if __name__ == '__main__':
    create_presentation()
